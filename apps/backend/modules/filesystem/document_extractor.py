import os
import io
import logging
from typing import List, Dict, Any, Optional

logger = logging.getLogger("JARVIS.DocumentExtractor")

class ExtractionResult:
    def __init__(
        self,
        text: str = "",
        tables: List[List[List[Any]]] = None,
        metadata: Dict[str, Any] = None,
        method: str = "fallback_metadata",
        success: bool = False,
        error: Optional[str] = None
    ):
        self.text = text
        self.tables = tables or []
        self.metadata = metadata or {}
        self.method = method
        self.success = success
        self.error = error

    def to_dict(self) -> Dict[str, Any]:
        return {
            "text": self.text,
            "tables": self.tables,
            "metadata": self.metadata,
            "method": self.method,
            "success": self.success,
            "error": self.error
        }


class DocumentExtractor:
    def __init__(self):
        self._ocr_service = None

    def _get_translation_service(self):
        try:
            from container import ServiceContainer
            container = ServiceContainer.instance()
            if container:
                return container.get_or_none("translation_service")
        except Exception as e:
            logger.debug(f"Could not load translation_service from container: {e}")
        
        try:
            from modules.language.translation_service import TranslationService
            return TranslationService()
        except Exception as e:
            logger.error(f"Failed to instantiate TranslationService fallback: {e}")
        return None

    def _get_preferred_language(self) -> list[str]:
        try:
            from container import ServiceContainer
            container = ServiceContainer.instance()
            if container:
                memory = container.get_or_none("memory")
                if memory:
                    pref = memory.get_preference("preferred_language")
                    if pref:
                        return [pref, "en"]
        except Exception as e:
            logger.debug(f"Failed to fetch preferred language from memory: {e}")
        return ["en"]

    def _get_ocr_service(self):
        if self._ocr_service is not None:
            return self._ocr_service
        try:
            from container import ServiceContainer
            container = ServiceContainer.instance()
            if container:
                vision_manager = container.get_or_none("vision_manager")
                if vision_manager and hasattr(vision_manager, "ocr_service"):
                    self._ocr_service = vision_manager.ocr_service
                    return self._ocr_service
        except Exception as e:
            logger.debug(f"Could not load OCRService from container: {e}")
            
        try:
            from modules.vision.ocr_service import OCRService
            self._ocr_service = OCRService()
        except Exception as e:
            logger.error(f"Failed to instantiate OCRService fallback: {e}")
            
        return self._ocr_service

    def extract(self, file_path: str, **kwargs) -> ExtractionResult:
        """Route file by extension and magic bytes, return structured result."""
        if not os.path.exists(file_path):
            return ExtractionResult(
                text="",
                success=False,
                error=f"File not found: {file_path}",
                method="fallback_metadata"
            )
            
        ext = self.detect_file_type(file_path)
        logger.info(f"Extracting content from '{file_path}' (detected type: {ext})")
        
        languages = kwargs.get("languages") or kwargs.get("languages_list")
        if not languages:
            languages = self._get_preferred_language()
        
        if ext == ".pdf":
            res = self._extract_pdf(file_path, languages=languages)
        elif ext == ".docx":
            ocr_embedded = kwargs.get("ocr_embedded_images", False)
            res = self._extract_docx(file_path, ocr_embedded_images=ocr_embedded, languages=languages)
        elif ext == ".xlsx":
            res = self._extract_xlsx(file_path)
        elif ext == ".pptx":
            res = self._extract_pptx(file_path)
        elif ext == ".csv":
            res = self._extract_csv(file_path)
        elif ext == ".json":
            res = self._extract_json(file_path)
        elif ext in (".yaml", ".yml"):
            res = self._extract_yaml(file_path)
        elif ext == ".xml":
            res = self._extract_xml(file_path)
        elif ext in (".html", ".htm"):
            res = self._extract_html(file_path)
        elif ext in (".txt", ".py", ".md", ".sh", ".bat", ".ini", ".cfg", ".conf", ".sql"):
            res = self._extract_text_file(file_path)
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp"):
            res = self._extract_image(file_path, languages=languages)
        else:
            # Fallback metadata representation
            try:
                size = os.path.getsize(file_path)
                filename = os.path.basename(file_path)
                res = ExtractionResult(
                    text=f"Binary file name: {filename}, Extension: {ext}, Size: {size} bytes.",
                    metadata={"file_path": file_path, "filename": filename, "extension": ext, "size_bytes": size},
                    method="fallback_metadata",
                    success=True
                )
            except Exception as e:
                res = ExtractionResult(
                    text="",
                    success=False,
                    error=str(e),
                    method="fallback_metadata"
                )

        translate_to = kwargs.get("translate_to")
        if translate_to and res.success and res.text:
            trans_service = self._get_translation_service()
            if trans_service:
                trans_res = trans_service.translate(res.text, target_lang=translate_to)
                res.text = trans_res.translated_text
                if res.metadata is None:
                    res.metadata = {}
                res.metadata["translation"] = {
                    "source_lang": trans_res.source_lang,
                    "target_lang": trans_res.target_lang,
                    "engine_used": trans_res.engine_used
                }
        return res

    def detect_file_type(self, file_path: str) -> str:
        """
        Detects file type by extension first.
        If extension is missing or not recognized, checks magic bytes.
        """
        ext = os.path.splitext(file_path)[1].lower()
        if ext in (".pdf", ".docx", ".xlsx", ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".zip", ".pptx", ".csv", ".json", ".yaml", ".yml", ".xml", ".html", ".htm"):
            return ext
            
        # Check magic bytes
        try:
            with open(file_path, "rb") as f:
                header = f.read(4)
            if header.startswith(b'\x25\x50\x44\x46'): # %PDF
                return ".pdf"
            elif header.startswith(b'\x89\x50\x4e\x47'): # PNG
                return ".png"
            elif header.startswith(b'\xff\xd8\xff'): # JPEG
                return ".jpeg"
            elif header.startswith(b'\x47\x49\x46\x38'): # GIF
                return ".gif"
            elif header.startswith(b'\x42\x4d'): # BMP
                return ".bmp"
            elif header.startswith(b'\x50\x4b\x03\x04'): # PK (ZIP / DOCX / XLSX / PPTX)
                import zipfile
                if zipfile.is_zipfile(file_path):
                    with zipfile.ZipFile(file_path) as z:
                        namelist = z.namelist()
                        if any("word/document.xml" in name for name in namelist):
                            return ".docx"
                        elif any("xl/workbook.xml" in name for name in namelist):
                            return ".xlsx"
                        elif any("ppt/presentation.xml" in name for name in namelist):
                            return ".pptx"
                    return ".zip"
        except Exception as e:
            logger.warning(f"Error detecting file type via magic bytes for {file_path}: {e}")
            
        return ext

    def _extract_pdf(self, file_path: str, languages: list[str] = None) -> ExtractionResult:
        try:
            import pypdf
            import pdfplumber
            from PIL import Image
        except ImportError as err:
            logger.error(f"Required PDF dependencies missing: {err}")
            return ExtractionResult(
                text="",
                success=False,
                error=f"Required PDF dependencies missing: {err}",
                method="fallback_metadata"
            )

        text_pages = []
        tables = []
        metadata = {
            "file_path": file_path,
            "filename": os.path.basename(file_path),
            "page_count": 0,
            "pages": []
        }
        method = "native"
        
        try:
            reader = pypdf.PdfReader(file_path)
            page_count = len(reader.pages)
            metadata["page_count"] = page_count
            
            # Auto-detect language if not explicitly provided
            langs_to_use = languages
            if not langs_to_use:
                try:
                    sample_text = ""
                    for p in reader.pages:
                        sample_text += (p.extract_text() or "")
                        if len(sample_text) > 200:
                            break
                    if sample_text:
                        from modules.language.language_detector import detect_language
                        lang_res = detect_language(sample_text)
                        if lang_res.code in ["hi", "gu"] and lang_res.confidence >= 0.5:
                            langs_to_use = [lang_res.code]
                            logger.info(f"Auto-detected language for PDF: {lang_res.code} (confidence: {lang_res.confidence})")
                except Exception as detect_err:
                    logger.debug(f"Failed to auto-detect PDF language: {detect_err}")
            
            if not langs_to_use:
                langs_to_use = ["en"]

            for idx, page in enumerate(reader.pages):
                page_num = idx + 1
                page_text = page.extract_text() or ""
                page_method = "native"
                
                # If page text is very short/empty, try OCR on page images
                if len(page_text.strip()) < 15:
                    ocr_service = self._get_ocr_service()
                    if ocr_service:
                        ocr_texts = []
                        # page.images is a list of image objects in pypdf
                        for img_idx, img_obj in enumerate(page.images):
                            try:
                                img = Image.open(io.BytesIO(img_obj.data))
                                ocr_res = ocr_service.extract_text(img, languages=langs_to_use)
                                if ocr_res and ocr_res.strip():
                                    ocr_texts.append(ocr_res)
                            except Exception as img_err:
                                logger.warning(f"Error OCR-ing page {page_num} image {img_idx}: {img_err}")
                        if ocr_texts:
                            page_text = "\n".join(ocr_texts)
                            page_method = "ocr"
                            method = "ocr"
                
                text_pages.append(page_text)
                metadata["pages"].append({
                    "page_number": page_num,
                    "length": len(page_text),
                    "method": page_method
                })
            
            # Now let's try to extract tables using pdfplumber
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        extracted_tables = page.extract_tables()
                        for table in extracted_tables:
                            if table:
                                tables.append(table)
            except Exception as table_err:
                logger.warning(f"Error extracting tables via pdfplumber: {table_err}")

            full_text = ""
            for idx, page_text in enumerate(text_pages):
                page_num = idx + 1
                if page_text.strip():
                    full_text += f"\n--- Page {page_num} ---\n{page_text}\n"

            metadata["page_texts"] = text_pages

            return ExtractionResult(
                text=full_text.strip(),
                tables=tables,
                metadata=metadata,
                method=method,
                success=True
            )
        except Exception as e:
            logger.error(f"Error reading PDF {file_path}: {e}")
            return ExtractionResult(
                text="",
                metadata={"file_path": file_path, "filename": os.path.basename(file_path)},
                method="fallback_metadata",
                success=False,
                error=str(e)
            )

    def _extract_docx(self, file_path: str, ocr_embedded_images: bool = False, languages: list[str] = None) -> ExtractionResult:
        try:
            import docx
            from PIL import Image
        except ImportError as err:
            logger.error(f"Required DOCX dependencies missing: {err}")
            return ExtractionResult(
                text="",
                success=False,
                error=f"Required DOCX dependencies missing: {err}",
                method="fallback_metadata"
            )

        paragraphs_text = []
        tables = []
        metadata = {
            "file_path": file_path,
            "filename": os.path.basename(file_path),
            "paragraph_count": 0,
            "table_count": 0
        }
        method = "native"
        
        try:
            doc = docx.Document(file_path)
            
            # Auto-detect language if not explicitly provided
            langs_to_use = languages
            if not langs_to_use:
                try:
                    sample_text = ""
                    for p in doc.paragraphs:
                        sample_text += p.text
                        if len(sample_text) > 200:
                            break
                    if sample_text:
                        from modules.language.language_detector import detect_language
                        lang_res = detect_language(sample_text)
                        if lang_res.code in ["hi", "gu"] and lang_res.confidence >= 0.5:
                            langs_to_use = [lang_res.code]
                            logger.info(f"Auto-detected language for DOCX: {lang_res.code} (confidence: {lang_res.confidence})")
                except Exception as detect_err:
                    logger.debug(f"Failed to auto-detect DOCX language: {detect_err}")
            
            if not langs_to_use:
                langs_to_use = ["en"]

            for p in doc.paragraphs:
                if p.text.strip():
                    paragraphs_text.append(p.text)
            
            metadata["paragraph_count"] = len(doc.paragraphs)
            metadata["table_count"] = len(doc.tables)
            
            # Extract tables
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    tables.append(table_data)
                    
            # Optionally OCR embedded images
            if ocr_embedded_images:
                ocr_service = self._get_ocr_service()
                if ocr_service:
                    for rel in doc.part.rels.values():
                        if "image" in rel.target_ref:
                            try:
                                img_data = rel.target_part.blob
                                img = Image.open(io.BytesIO(img_data))
                                ocr_res = ocr_service.extract_text(img, languages=langs_to_use)
                                if ocr_res and ocr_res.strip():
                                    paragraphs_text.append(f"\n[Embedded Image OCR]\n{ocr_res}\n")
                                    method = "ocr"
                            except Exception as img_err:
                                logger.warning(f"Error OCR-ing DOCX embedded image: {img_err}")
                                
            full_text = "\n".join(paragraphs_text)
            if tables:
                full_text += "\n\n--- Extracted Tables ---\n"
                for idx, t in enumerate(tables):
                    full_text += f"\nTable {idx + 1}:\n"
                    for row in t:
                        full_text += " | ".join(row) + "\n"
                        
            return ExtractionResult(
                text=full_text.strip(),
                tables=tables,
                metadata=metadata,
                method=method,
                success=True
            )
        except Exception as e:
            logger.error(f"Error reading DOCX {file_path}: {e}")
            return ExtractionResult(
                text="",
                metadata={"file_path": file_path, "filename": os.path.basename(file_path)},
                method="fallback_metadata",
                success=False,
                error=str(e)
            )

    def _extract_xlsx(self, file_path: str) -> ExtractionResult:
        try:
            import pandas as pd
            import openpyxl
        except ImportError as err:
            logger.error(f"Required XLSX dependencies missing: {err}")
            return ExtractionResult(
                text="",
                success=False,
                error=f"Required XLSX dependencies missing: {err}",
                method="fallback_metadata"
            )

        tables = []
        metadata = {
            "file_path": file_path,
            "filename": os.path.basename(file_path),
            "sheet_names": []
        }
        method = "native"
        
        try:
            with pd.ExcelFile(file_path) as xlsx:
                sheet_names = xlsx.sheet_names
                metadata["sheet_names"] = sheet_names
                
                sheet_texts = {}
                full_text = ""
                for sheet_name in sheet_names:
                    df = pd.read_excel(xlsx, sheet_name=sheet_name)
                    df = df.fillna("")
                    
                    sheet_rows = [df.columns.tolist()] + df.values.tolist()
                    tables.append(sheet_rows)
                    
                    sheet_text = f"\n--- Sheet: {sheet_name} ---\n"
                    headers = " | ".join(str(h) for h in df.columns)
                    sheet_text += headers + "\n"
                    sheet_text += "-|-".join("-" * len(str(h)) for h in df.columns) + "\n"
                    for _, row in df.iterrows():
                        sheet_text += " | ".join(str(val) for val in row.values) + "\n"
                    
                    sheet_texts[sheet_name] = sheet_text.strip()
                    full_text += sheet_text
                    
                metadata["sheet_texts"] = sheet_texts
                
                return ExtractionResult(
                    text=full_text.strip(),
                    tables=tables,
                    metadata=metadata,
                    method=method,
                    success=True
                )
        except Exception as e:
            logger.error(f"Error reading XLSX {file_path}: {e}")
            return ExtractionResult(
                text="",
                metadata={"file_path": file_path, "filename": os.path.basename(file_path)},
                method="fallback_metadata",
                success=False,
                error=str(e)
            )

    def _extract_image(self, file_path: str, languages: list[str] = None) -> ExtractionResult:
        try:
            from PIL import Image
        except ImportError as err:
            logger.error(f"Required PIL dependency missing: {err}")
            return ExtractionResult(
                text="",
                success=False,
                error=f"Required PIL dependency missing: {err}",
                method="fallback_metadata"
            )

        metadata = {
            "file_path": file_path,
            "filename": os.path.basename(file_path)
        }
        
        try:
            img = Image.open(file_path)
            metadata["size_pixels"] = img.size
            metadata["format"] = img.format
            
            ocr_service = self._get_ocr_service()
            if ocr_service:
                text = ocr_service.extract_text(img, languages=languages)
                method = "ocr"
            else:
                text = ""
                method = "fallback_metadata"
                
            return ExtractionResult(
                text=text.strip(),
                metadata=metadata,
                method=method,
                success=True
            )
        except Exception as e:
            logger.error(f"Error OCR-ing Image {file_path}: {e}")
            return ExtractionResult(
                text="",
                metadata={"file_path": file_path, "filename": os.path.basename(file_path)},
                method="fallback_metadata",
                success=False,
                error=str(e)
            )

    def _extract_pptx(self, file_path: str) -> ExtractionResult:
        try:
            from pptx import Presentation
        except ImportError as err:
            logger.error(f"Required PPTX dependencies missing: {err}")
            return ExtractionResult(
                text="",
                success=False,
                error=f"Required PPTX dependencies missing: {err}",
                method="fallback_metadata"
            )
            
        metadata = {
            "file_path": file_path,
            "filename": os.path.basename(file_path),
            "slide_count": 0
        }
        
        try:
            prs = Presentation(file_path)
            slide_count = len(prs.slides)
            metadata["slide_count"] = slide_count
            
            full_text = ""
            for idx, slide in enumerate(prs.slides):
                slide_num = idx + 1
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    full_text += f"\n--- Slide {slide_num} ---\n"
                    full_text += "\n".join(slide_text) + "\n"
                    
            return ExtractionResult(
                text=full_text.strip(),
                metadata=metadata,
                method="native",
                success=True
            )
        except Exception as e:
            logger.error(f"Error reading PPTX {file_path}: {e}")
            return ExtractionResult(
                text="",
                metadata={"file_path": file_path, "filename": os.path.basename(file_path)},
                method="fallback_metadata",
                success=False,
                error=str(e)
            )

    def _extract_csv(self, file_path: str) -> ExtractionResult:
        try:
            import pandas as pd
        except ImportError as err:
            logger.error(f"Required pandas dependency missing: {err}")
            return ExtractionResult(
                text="",
                success=False,
                error=f"Required pandas dependency missing: {err}",
                method="fallback_metadata"
            )
            
        metadata = {
            "file_path": file_path,
            "filename": os.path.basename(file_path)
        }
        
        try:
            df = pd.read_csv(file_path)
            df = df.fillna("")
            
            rows = [df.columns.tolist()] + df.values.tolist()
            
            headers = " | ".join(str(h) for h in df.columns)
            markdown_text = headers + "\n"
            markdown_text += "-|-".join("-" * len(str(h)) for h in df.columns) + "\n"
            for _, row in df.iterrows():
                markdown_text += " | ".join(str(val) for val in row.values) + "\n"
                
            metadata["row_count"] = len(df)
            metadata["column_count"] = len(df.columns)
            
            return ExtractionResult(
                text=markdown_text.strip(),
                tables=[rows],
                metadata=metadata,
                method="native",
                success=True
            )
        except Exception as e:
            logger.error(f"Error reading CSV {file_path}: {e}")
            return ExtractionResult(
                text="",
                metadata={"file_path": file_path, "filename": os.path.basename(file_path)},
                method="fallback_metadata",
                success=False,
                error=str(e)
            )

    def _extract_json(self, file_path: str) -> ExtractionResult:
        import json
        metadata = {
            "file_path": file_path,
            "filename": os.path.basename(file_path)
        }
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                data = json.load(f)
            pretty_text = json.dumps(data, indent=2)
            return ExtractionResult(
                text=pretty_text,
                metadata=metadata,
                method="native",
                success=True
            )
        except Exception as e:
            logger.error(f"Error reading JSON {file_path}: {e}")
            return ExtractionResult(
                text="",
                metadata={"file_path": file_path, "filename": os.path.basename(file_path)},
                method="fallback_metadata",
                success=False,
                error=str(e)
            )

    def _extract_yaml(self, file_path: str) -> ExtractionResult:
        metadata = {
            "file_path": file_path,
            "filename": os.path.basename(file_path)
        }
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            return ExtractionResult(
                text=content.strip(),
                metadata=metadata,
                method="native",
                success=True
            )
        except Exception as e:
            logger.error(f"Error reading YAML {file_path}: {e}")
            return ExtractionResult(
                text="",
                metadata={"file_path": file_path, "filename": os.path.basename(file_path)},
                method="fallback_metadata",
                success=False,
                error=str(e)
            )

    def _extract_xml(self, file_path: str) -> ExtractionResult:
        metadata = {
            "file_path": file_path,
            "filename": os.path.basename(file_path)
        }
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            return ExtractionResult(
                text=content.strip(),
                metadata=metadata,
                method="native",
                success=True
            )
        except Exception as e:
            logger.error(f"Error reading XML {file_path}: {e}")
            return ExtractionResult(
                text="",
                metadata={"file_path": file_path, "filename": os.path.basename(file_path)},
                method="fallback_metadata",
                success=False,
                error=str(e)
            )

    def _extract_html(self, file_path: str) -> ExtractionResult:
        import re
        metadata = {
            "file_path": file_path,
            "filename": os.path.basename(file_path)
        }
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                html = f.read()
                
            clean_html = re.sub(r'<(script|style).*?>.*?</\1>', '', html, flags=re.DOTALL | re.IGNORECASE)
            clean_text = re.sub(r'<.*?>', '', clean_html)
            clean_text = clean_text.replace("&nbsp;", " ").replace("&lt;", "<").replace("&gt;", ">").replace("&amp;", "&")
            clean_text = re.sub(r'\n\s*\n', '\n\n', clean_text)
            
            return ExtractionResult(
                text=clean_text.strip(),
                metadata=metadata,
                method="native",
                success=True
            )
        except Exception as e:
            logger.error(f"Error reading HTML {file_path}: {e}")
            return ExtractionResult(
                text="",
                metadata={"file_path": file_path, "filename": os.path.basename(file_path)},
                method="fallback_metadata",
                success=False,
                error=str(e)
            )

    def _extract_text_file(self, file_path: str) -> ExtractionResult:
        metadata = {
            "file_path": file_path,
            "filename": os.path.basename(file_path)
        }
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            return ExtractionResult(
                text=content,
                metadata=metadata,
                method="native",
                success=True
            )
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}")
            return ExtractionResult(
                text="",
                metadata={"file_path": file_path, "filename": os.path.basename(file_path)},
                method="fallback_metadata",
                success=False,
                error=str(e)
            )
